#!/usr/bin/env python3
"""
Simple Document Text Extractor
A lightweight and efficient text extractor for common document formats.

Author: BlackBoxAI
License: MIT
"""

import os
import sys
import argparse
from pathlib import Path
from typing import Optional, Dict, List, Union
from concurrent.futures import ThreadPoolExecutor, as_completed
import json
from datetime import datetime

# Core dependencies - using more efficient libraries
try:
    import fitz  # PyMuPDF - faster PDF extraction
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class SimpleDocumentExtractor:
    """Efficient document text extractor with minimal dependencies."""
    
    SUPPORTED_FORMATS = {
        '.pdf': 'extract_pdf',
        '.docx': 'extract_docx',
        '.pptx': 'extract_pptx',
        '.txt': 'extract_text',
        '.md': 'extract_text',
        '.csv': 'extract_text',
        '.log': 'extract_text',
        '.json': 'extract_text',
        '.xml': 'extract_text',
        '.html': 'extract_text',
        '.htm': 'extract_text'
    }
    
    def __init__(self, max_workers: int = 4):
        """
        Initialize the extractor.
        
        Args:
            max_workers: Maximum number of threads for parallel processing
        """
        self.max_workers = max_workers
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check available dependencies and update supported formats."""
        if not PDF_AVAILABLE:
            for ext in ['.pdf']:
                if ext in self.SUPPORTED_FORMATS:
                    del self.SUPPORTED_FORMATS[ext]
        
        if not DOCX_AVAILABLE:
            for ext in ['.docx']:
                if ext in self.SUPPORTED_FORMATS:
                    del self.SUPPORTED_FORMATS[ext]
        
        if not PPTX_AVAILABLE:
            for ext in ['.pptx']:
                if ext in self.SUPPORTED_FORMATS:
                    del self.SUPPORTED_FORMATS[ext]
    
    def extract(self, path: Union[str, Path]) -> Dict[str, str]:
        """
        Extract text from a file or directory.
        
        Args:
            path: File or directory path
            
        Returns:
            Dictionary with file paths as keys and extracted text as values
        """
        path = Path(path)
        
        if path.is_file():
            return self._extract_single_file(path)
        elif path.is_dir():
            return self._extract_directory(path)
        else:
            raise ValueError(f"Path does not exist: {path}")
    
    def _extract_single_file(self, file_path: Path) -> Dict[str, str]:
        """Extract text from a single file."""
        ext = file_path.suffix.lower()
        
        if ext not in self.SUPPORTED_FORMATS:
            return {str(file_path): f"Unsupported format: {ext}"}
        
        method_name = self.SUPPORTED_FORMATS[ext]
        method = getattr(self, method_name)
        
        try:
            text = method(file_path)
            return {str(file_path): text}
        except Exception as e:
            return {str(file_path): f"Error: {str(e)}"}
    
    def _extract_directory(self, dir_path: Path) -> Dict[str, str]:
        """Extract text from all supported files in a directory."""
        results = {}
        files = []
        
        # Collect all supported files
        for ext in self.SUPPORTED_FORMATS:
            files.extend(dir_path.rglob(f"*{ext}"))
        
        # Process files in parallel
        if len(files) > 1 and self.max_workers > 1:
            with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
                future_to_file = {
                    executor.submit(self._extract_single_file, f): f 
                    for f in files
                }
                
                for future in as_completed(future_to_file):
                    result = future.result()
                    results.update(result)
        else:
            # Process sequentially for single file or single thread
            for file_path in files:
                result = self._extract_single_file(file_path)
                results.update(result)
        
        return results
    
    def extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF using PyMuPDF (faster than PyPDF2/pdfplumber)."""
        if not PDF_AVAILABLE:
            raise ImportError("PyMuPDF not installed. Install with: pip install pymupdf")
        
        text_parts = []
        
        with fitz.open(file_path) as pdf:
            for page_num, page in enumerate(pdf, 1):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"[Page {page_num}]\n{text}")
        
        return '\n'.join(text_parts)
    
    def extract_docx(self, file_path: Path) -> str:
        """Extract text from Word document."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
        
        doc = Document(file_path)
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables (simplified)
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    
    def extract_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[Slide {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if len(slide_texts) > 1:  # Has content beyond slide number
                text_parts.extend(slide_texts)
        
        return '\n'.join(text_parts)
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from plain text files with automatic encoding detection."""
        # Try common encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        # If all fail, try with errors='ignore'
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()


class BatchExtractor:
    """Optimized batch extraction with progress tracking."""
    
    def __init__(self, extractor: SimpleDocumentExtractor):
        self.extractor = extractor
        self.stats = {
            'processed': 0,
            'failed': 0,
            'total_size': 0
        }
    
    def extract_batch(self, paths: List[Union[str, Path]], 
                     output_format: str = 'text') -> Union[str, Dict]:
        """
        Extract text from multiple files efficiently.
        
        Args:
            paths: List of file or directory paths
            output_format: 'text', 'json', or 'dict'
            
        Returns:
            Extracted text in specified format
        """
        all_results = {}
        
        for path in paths:
            results = self.extractor.extract(path)
            all_results.update(results)
            
            # Update stats
            for file_path, content in results.items():
                if not content.startswith("Error:"):
                    self.stats['processed'] += 1
                    self.stats['total_size'] += len(content)
                else:
                    self.stats['failed'] += 1
        
        # Format output
        if output_format == 'dict':
            return all_results
        elif output_format == 'json':
            return json.dumps({
                'timestamp': datetime.now().isoformat(),
                'stats': self.stats,
                'results': all_results
            }, indent=2)
        else:  # text format
            output_parts = []
            for file_path, content in all_results.items():
                output_parts.append(f"\n{'='*60}")
                output_parts.append(f"FILE: {file_path}")
                output_parts.append(f"{'='*60}")
                output_parts.append(content)
            return '\n'.join(output_parts)


def quick_extract(path: Union[str, Path], save_to: Optional[str] = None) -> str:
    """
    Quick extraction function for simple use cases.
    
    Args:
        path: File or directory to extract from
        save_to: Optional output file path
        
    Returns:
        Extracted text
    """
    extractor = SimpleDocumentExtractor()
    results = extractor.extract(path)
    
    # Combine all text
    if len(results) == 1:
        text = list(results.values())[0]
    else:
        text = '\n\n'.join(f"[{k}]\n{v}" for k, v in results.items())
    
    # Save if requested
    if save_to:
        Path(save_to).write_text(text, encoding='utf-8')
        print(f"Saved to: {save_to}")
    
    return text


def main():
    """Command-line interface."""
    parser = argparse.ArgumentParser(
        description='Simple and efficient document text extractor',
        epilog="""
Examples:
  # Extract from single file
  python simple_document_extractor.py document.pdf
  
  # Extract from directory
  python simple_document_extractor.py /path/to/documents/
  
  # Save output
  python simple_document_extractor.py document.pdf -o output.txt
  
  # JSON output
  python simple_document_extractor.py /path/to/docs/ --json
  
  # Parallel processing
  python simple_document_extractor.py /path/to/docs/ -w 8
        """
    )
    
    parser.add_argument('paths', nargs='+', help='Files or directories to extract from')
    parser.add_argument('-o', '--output', help='Output file path')
    parser.add_argument('--json', action='store_true', help='Output as JSON')
    parser.add_argument('-w', '--workers', type=int, default=4,
                       help='Number of parallel workers (default: 4)')
    parser.add_argument('--check', action='store_true', 
                       help='Check available dependencies')
    
    args = parser.parse_args()
    
    if args.check:
        print("Dependency Status:")
        print(f"  PDF support (PyMuPDF): {'✓' if PDF_AVAILABLE else '✗'}")
        print(f"  Word support (python-docx): {'✓' if DOCX_AVAILABLE else '✗'}")
        print(f"  PowerPoint support (python-pptx): {'✓' if PPTX_AVAILABLE else '✗'}")
        
        if not PDF_AVAILABLE:
            print("\nInstall PyMuPDF for PDF support: pip install pymupdf")
        if not DOCX_AVAILABLE:
            print("Install python-docx for Word support: pip install python-docx")
        if not PPTX_AVAILABLE:
            print("Install python-pptx for PowerPoint support: pip install python-pptx")
        return
    
    try:
        # Create extractor
        extractor = SimpleDocumentExtractor(max_workers=args.workers)
        batch = BatchExtractor(extractor)
        
        # Process files
        output_format = 'json' if args.json else 'text'
        result = batch.extract_batch(args.paths, output_format=output_format)
        
        # Output results
        if args.output:
            Path(args.output).write_text(result, encoding='utf-8')
            print(f"✓ Saved to: {args.output}")
        else:
            print(result)
        
        # Print stats
        print(f"\n--- Statistics ---")
        print(f"Files processed: {batch.stats['processed']}")
        if batch.stats['failed'] > 0:
            print(f"Files failed: {batch.stats['failed']}")
        print(f"Total characters: {batch.stats['total_size']:,}")
        
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
